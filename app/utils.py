import os
import re
import json
import csv
from xml.etree import ElementTree as ET
import chardet
import pytesseract
from pdf2image import convert_from_path
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load as odf_load
from pdf_handler import estrai_blocchi_testuali, rileva_tipo_pdf
from pdf_renderer import crea_pdf_bilingue, crea_html_bilingue

LINGUE_SUPPORTATE = {
    "ita": "Italiano",
    "eng": "Inglese",
    "fra": "Francese",
    "spa": "Spagnolo",
    "deu": "Tedesco",
    "por": "Portoghese",
    "rus": "Russo",
    "nld": "Olandese",
    "zho": "Cinese (semplificato)",
    "jpn": "Giapponese",
    "kor": "Coreano",
    "ara": "Arabo",
    "hin": "Hindi",
    "tur": "Turco",
    "pol": "Polacco",
    "swe": "Svedese",
}

LINGUE_OCR = {
    "ita": "Italiano",
    "eng": "Inglese",
    "fra": "Francese",
    "spa": "Spagnolo",
    "deu": "Tedesco",
    "por": "Portoghese",
    "rus": "Russo",
    "nld": "Olandese",
}

FORMATI_SUPPORTATI = {
    ".txt":  "Testo",
    ".pdf":  "PDF (OCR)",
    ".docx": "Word",
    ".xlsx": "Excel",
    ".xls":  "Excel (vecchio)",
    ".csv":  "CSV",
    ".odt":  "LibreOffice Writer",
    ".pptx": "PowerPoint",
    ".srt":  "Sottotitoli SRT",
    ".vtt":  "Sottotitoli VTT",
    ".html": "HTML",
    ".htm":  "HTML",
    ".md":   "Markdown",
    ".json": "JSON",
    ".xml":  "XML",
}

def _decodifica(filepath):
    with open(filepath, "rb") as f:
        raw = f.read()
    risultato = chardet.detect(raw)
    encoding = risultato.get("encoding") or "utf-8"
    try:
        return raw.decode(encoding)
    except:
        return raw.decode("utf-8", errors="replace")

def _scrivi_testo(testo, original_path, output_dir, suffisso):
    os.makedirs(output_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(original_path))[0]
    out_path = os.path.join(output_dir, f"{base}_{suffisso}.txt")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(testo)
    return out_path

# --- PARSER SOTTOTITOLI ---

def _analizza_srt(contenuto):
    blocchi = re.split(r"\n\s*\n", contenuto.strip())
    struttura = []
    testi = []
    for blocco in blocchi:
        righe = blocco.strip().split("\n")
        if len(righe) < 3:
            if righe and righe[0].strip():
                struttura.append(("", "", righe[0].strip()))
                testi.append(righe[0].strip())
            continue
        numero = righe[0].strip()
        timestamp = righe[1].strip()
        testo = "\n".join(righe[2:])
        struttura.append((numero, timestamp, testo))
        testi.append(testo)
    return struttura, testi

def _ricostruisci_srt(struttura, testi_tradotti):
    risultato = []
    for i, (numero, timestamp, _) in enumerate(struttura):
        blocco = [numero, timestamp, testi_tradotti[i] if i < len(testi_tradotti) else ""]
        risultato.append("\n".join(blocco))
    return "\n\n".join(risultato)

def _analizza_vtt(contenuto):
    righe = contenuto.strip().split("\n")
    struttura = []
    testi = []
    i = 0
    intestazione = ""
    if righe and righe[0].strip() == "WEBVTT":
        intestazione = righe[0]
        i = 1
    while i < len(righe):
        if righe[i].strip() == "":
            i += 1
            continue
        if "-->" in righe[i]:
            timestamp = righe[i].strip()
            i += 1
            testo_righe = []
            while i < len(righe) and righe[i].strip() != "":
                testo_righe.append(righe[i])
                i += 1
            testo = "\n".join(testo_righe)
            struttura.append((intestazione if i == 1 else "", timestamp, testo))
            testi.append(testo)
        else:
            i += 1
    return struttura, testi, intestazione

def _ricostruisci_vtt(struttura, testi_tradotti):
    risultato = ["WEBVTT"]
    for i, (_, timestamp, _) in enumerate(struttura):
        risultato.append("")
        risultato.append(timestamp)
        if i < len(testi_tradotti):
            risultato.append(testi_tradotti[i])
    return "\n".join(risultato)

# --- LETTURA FILE ---

def leggi_testo(filepath):
    return _decodifica(filepath)

def leggi_docx(filepath):
    doc = Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs)

def leggi_pdf(filepath, lingua_ocr="ita+eng"):
    tipo = rileva_tipo_pdf(filepath)
    if tipo == "scansionato":
        immagini = convert_from_path(filepath)
        testo = ""
        for img in immagini:
            testo += pytesseract.image_to_string(img, lang=lingua_ocr) + "\n"
        return testo.strip()
    pagine = estrai_blocchi_testuali(filepath, lingua_ocr)
    return "\n".join(
        blocco["testo"] for pagina in pagine for blocco in pagina
    )

def leggi_xlsx(filepath):
    wb = load_workbook(filepath, data_only=True)
    parti = []
    for ws in wb.worksheets:
        for riga in ws.iter_rows():
            for cella in riga:
                if cella.value is not None and isinstance(cella.value, str) and cella.value.strip():
                    parti.append(cella.value.strip())
    return "\n".join(parti)

def leggi_xls(filepath):
    import xlrd
    wb = xlrd.open_workbook(filepath)
    parti = []
    for i in range(wb.nsheets):
        ws = wb.sheet_by_index(i)
        for r in range(ws.nrows):
            for c in range(ws.ncols):
                val = ws.cell_value(r, c)
                if isinstance(val, str) and val.strip():
                    parti.append(val.strip())
    return "\n".join(parti)

def leggi_csv(filepath):
    contenuto = _decodifica(filepath)
    return "\n".join(
        cella.strip() for riga in csv.reader(contenuto.splitlines())
        for cella in riga if cella.strip()
    )

def leggi_odt(filepath):
    doc = odf_load(filepath)
    parti = []
    for par in doc.getElementsByType(text.P):
        parti.append(teletype.extractText(par))
    return "\n".join(parti)

def leggi_pptx(filepath):
    prs = Presentation(filepath)
    parti = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for par in shape.text_frame.paragraphs:
                    parti.append(par.text)
    return "\n".join(parti)

def leggi_srt(filepath):
    contenuto = _decodifica(filepath)
    struttura, testi = _analizza_srt(contenuto)
    return "\n".join(testi)

def leggi_vtt(filepath):
    contenuto = _decodifica(filepath)
    struttura, testi, _ = _analizza_vtt(contenuto)
    return "\n".join(testi)

def leggi_html(filepath):
    contenuto = _decodifica(filepath)
    pulito = re.sub(r"<[^>]+>", "", contenuto)
    return re.sub(r"\s+", " ", pulito).strip()

def leggi_md(filepath):
    return _decodifica(filepath)

def leggi_json(filepath):
    contenuto = _decodifica(filepath)
    dati = json.loads(contenuto)
    def estrai_testo(obj):
        if isinstance(obj, str):
            return obj
        elif isinstance(obj, dict):
            return "\n".join(estrai_testo(v) for v in obj.values())
        elif isinstance(obj, list):
            return "\n".join(estrai_testo(v) for v in obj)
        return ""
    return estrai_testo(dati)

def leggi_xml(filepath):
    tree = ET.parse(filepath)
    testo = re.sub(r"<[^>]+>", "", ET.tostring(tree.getroot(), encoding="unicode"))
    return re.sub(r"\s+", " ", testo).strip()

# --- DISPATCH LETTURA ---

def leggi_file(filepath, lingua_ocr="ita+eng"):
    est = os.path.splitext(filepath)[1].lower()
    if est == ".txt":
        return leggi_testo(filepath)
    elif est == ".docx":
        return leggi_docx(filepath)
    elif est == ".pdf":
        return leggi_pdf(filepath, lingua_ocr)
    elif est == ".xlsx":
        return leggi_xlsx(filepath)
    elif est == ".xls":
        return leggi_xls(filepath)
    elif est == ".csv":
        return leggi_csv(filepath)
    elif est == ".odt":
        return leggi_odt(filepath)
    elif est == ".pptx":
        return leggi_pptx(filepath)
    elif est == ".srt":
        return leggi_srt(filepath)
    elif est == ".vtt":
        return leggi_vtt(filepath)
    elif est in (".html", ".htm"):
        return leggi_html(filepath)
    elif est == ".md":
        return leggi_md(filepath)
    elif est == ".json":
        return leggi_json(filepath)
    elif est == ".xml":
        return leggi_xml(filepath)
    raise ValueError(f"Formato non supportato: {est}")

# --- SCRITTURA OUTPUT ---

def scrivi_traduzione(testo_tradotto, file_originale, directory_output, lingua_target,
                     modalita_pdf="bilingue"):
    est = os.path.splitext(file_originale)[1].lower()
    suffisso = f"tradotto_{lingua_target}"

    if est == ".pdf":
        tipo = rileva_tipo_pdf(file_originale)
        if tipo == "scansionato":
            return _scrivi_testo(testo_tradotto, file_originale, directory_output, suffisso)

        pdf_path = crea_pdf_bilingue(
            file_originale, testo_tradotto, directory_output, lingua_target,
            modalita="affiancato" if modalita_pdf == "bilingue" else "mono"
        )
        html_path = crea_html_bilingue(
            file_originale, testo_tradotto, directory_output, lingua_target
        )
        return pdf_path

    if est == ".docx":
        doc = Document(file_originale)
        paragrafi = [p for p in doc.paragraphs if p.text.strip()]
        righe = [r for r in testo_tradotto.split("\n") if r.strip()]
        for i, par in enumerate(paragrafi):
            if i < len(righe):
                par.text = righe[i]
        base = os.path.splitext(os.path.basename(file_originale))[0]
        out_path = os.path.join(directory_output, f"{base}_{suffisso}.docx")
        doc.save(out_path)
        return out_path

    if est == ".xlsx":
        wb = load_workbook(file_originale)
        celle = []
        for ws in wb.worksheets:
            for riga in ws.iter_rows():
                for cella in riga:
                    if cella.value is not None and isinstance(cella.value, str) and cella.value.strip():
                        celle.append(cella)
        righe = [r for r in testo_tradotto.split("\n") if r.strip()]
        for i, cella in enumerate(celle):
            if i < len(righe):
                cella.value = righe[i]
        base = os.path.splitext(os.path.basename(file_originale))[0]
        out_path = os.path.join(directory_output, f"{base}_{suffisso}.xlsx")
        wb.save(out_path)
        return out_path

    if est == ".srt":
        contenuto = _decodifica(file_originale)
        struttura, testi_originali = _analizza_srt(contenuto)
        righe = [r for r in testo_tradotto.split("\n") if r.strip()]
        if len(righe) >= len(testi_originali):
            testi_tradotti = righe[:len(testi_originali)]
        else:
            testi_tradotti = righe + [""] * (len(testi_originali) - len(righe))
        srt_tradotto = _ricostruisci_srt(struttura, testi_tradotti)
        base = os.path.splitext(os.path.basename(file_originale))[0]
        out_path = os.path.join(directory_output, f"{base}_{suffisso}.srt")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(srt_tradotto)
        return out_path

    if est == ".vtt":
        contenuto = _decodifica(file_originale)
        struttura, testi_originali, _ = _analizza_vtt(contenuto)
        righe = [r for r in testo_tradotto.split("\n") if r.strip()]
        if len(righe) >= len(testi_originali):
            testi_tradotti = righe[:len(testi_originali)]
        else:
            testi_tradotti = righe + [""] * (len(testi_originali) - len(righe))
        vtt_tradotto = _ricostruisci_vtt(struttura, testi_tradotti)
        base = os.path.splitext(os.path.basename(file_originale))[0]
        out_path = os.path.join(directory_output, f"{base}_{suffisso}.vtt")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(vtt_tradotto)
        return out_path

    if est == ".csv":
        with open(file_originale, "r", encoding="utf-8") as f:
            righe_csv = list(csv.reader(f))
        celle_csv = [c for riga in righe_csv for c in riga if c.strip()]
        righe = [r for r in testo_tradotto.split("\n") if r.strip()]
        if len(righe) < len(celle_csv):
            righe += [""] * (len(celle_csv) - len(righe))
        indice = 0
        base = os.path.splitext(os.path.basename(file_originale))[0]
        out_path = os.path.join(directory_output, f"{base}_{suffisso}.csv")
        with open(out_path, "w", encoding="utf-8", newline="") as f:
            scrittore = csv.writer(f)
            for riga_orig in righe_csv:
                nuova = []
                for cella in riga_orig:
                    if cella.strip() and indice < len(righe):
                        nuova.append(righe[indice])
                        indice += 1
                    else:
                        nuova.append(cella)
                scrittore.writerow(nuova)
        return out_path

    return _scrivi_testo(testo_tradotto, file_originale, directory_output, suffisso)
